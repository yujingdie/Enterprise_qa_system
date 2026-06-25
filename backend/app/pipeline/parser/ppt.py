"""
PPT 解析器
每页幻灯片的文本提取
"""

from pptx import Presentation


def parse(file_path: str) -> list[dict]:
    """
    解析 PPT 文件，按页提取文本框文字

    返回格式:
    [
        {"page": 1, "text": "标题\n内容文字..."},
        {"page": 2, "text": "第二页内容..."},
    ]
    """
    prs = Presentation(file_path)
    pages = []

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)

        if texts:
            pages.append({
                "page": i + 1,
                "text": "\n".join(texts),
            })

    return pages
